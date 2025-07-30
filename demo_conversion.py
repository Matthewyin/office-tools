#!/usr/bin/env python
# -*- coding: utf-8 -*-

"""
转换功能演示脚本

演示Office转PDF工具的实际转换功能
"""

import tempfile
from pathlib import Path
from office2pdf import OfficeConverter


def create_sample_docx(file_path: Path):
    """创建示例Word文档"""
    try:
        from docx import Document
        
        doc = Document()
        doc.add_heading('Office转PDF工具演示', 0)
        
        doc.add_heading('功能特性', level=1)
        doc.add_paragraph('✅ 支持Word、Excel、PowerPoint转PDF')
        doc.add_paragraph('✅ 批量转换和递归处理')
        doc.add_paragraph('✅ 跨平台支持')
        doc.add_paragraph('✅ 详细的日志记录')
        
        doc.add_heading('技术规范', level=1)
        doc.add_paragraph('本项目严格遵循现代Python开发规范：')
        doc.add_paragraph('• 使用类型提示（Type Hints）')
        doc.add_paragraph('• 遵循PEP 8代码风格')
        doc.add_paragraph('• 完整的错误处理和日志记录')
        doc.add_paragraph('• 模块化设计和单一职责原则')
        
        doc.save(file_path)
        return True
    except Exception as e:
        print(f"创建Word文档失败: {e}")
        return False


def create_sample_xlsx(file_path: Path):
    """创建示例Excel文档"""
    try:
        from openpyxl import Workbook
        
        wb = Workbook()
        ws = wb.active
        ws.title = "转换统计"
        
        # 添加标题
        ws['A1'] = 'Office转PDF工具统计'
        ws['A1'].font = ws['A1'].font.copy(bold=True, size=16)
        
        # 添加数据
        data = [
            ['文件类型', '支持状态', '转换速度'],
            ['Word (.docx)', '✅ 支持', '快速'],
            ['Excel (.xlsx)', '✅ 支持', '快速'],
            ['PowerPoint (.pptx)', '✅ 支持', '快速'],
            ['PDF (.pdf)', '❌ 不需要', 'N/A'],
        ]
        
        for row_idx, row_data in enumerate(data, start=3):
            for col_idx, value in enumerate(row_data, start=1):
                ws.cell(row=row_idx, column=col_idx, value=value)
        
        # 设置列宽
        ws.column_dimensions['A'].width = 20
        ws.column_dimensions['B'].width = 15
        ws.column_dimensions['C'].width = 15
        
        wb.save(file_path)
        return True
    except Exception as e:
        print(f"创建Excel文档失败: {e}")
        return False


def create_sample_pptx(file_path: Path):
    """创建示例PowerPoint文档"""
    try:
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation()
        
        # 标题幻灯片
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "Office转PDF工具"
        subtitle.text = "现代Python开发规范实践\n作者: Matthew Yin"
        
        # 内容幻灯片
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = '主要功能'
        
        tf = body_shape.text_frame
        tf.text = '支持多种Office格式转换'
        
        p = tf.add_paragraph()
        p.text = '批量处理和递归转换'
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = '跨平台兼容性'
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = '完整的错误处理'
        p.level = 1
        
        prs.save(file_path)
        return True
    except Exception as e:
        print(f"创建PowerPoint文档失败: {e}")
        return False


def demo_conversion():
    """演示转换功能"""
    print("🚀 Office转PDF工具转换演示")
    print("=" * 50)
    
    with tempfile.TemporaryDirectory() as temp_dir:
        temp_path = Path(temp_dir)
        print(f"📁 工作目录: {temp_path}")
        
        # 创建示例文件
        print("\n📝 创建示例Office文件...")
        
        sample_files = []
        
        # Word文档
        docx_file = temp_path / "演示文档.docx"
        if create_sample_docx(docx_file):
            sample_files.append(docx_file)
            print(f"✅ 创建Word文档: {docx_file.name}")
        
        # Excel文档
        xlsx_file = temp_path / "统计表格.xlsx"
        if create_sample_xlsx(xlsx_file):
            sample_files.append(xlsx_file)
            print(f"✅ 创建Excel文档: {xlsx_file.name}")
        
        # PowerPoint文档
        pptx_file = temp_path / "演示幻灯片.pptx"
        if create_sample_pptx(pptx_file):
            sample_files.append(pptx_file)
            print(f"✅ 创建PowerPoint文档: {pptx_file.name}")
        
        if not sample_files:
            print("❌ 无法创建示例文件，演示结束")
            return
        
        # 创建转换器
        print(f"\n🔄 开始转换 {len(sample_files)} 个文件...")
        output_dir = temp_path / "pdf_output"
        converter = OfficeConverter(output_dir=str(output_dir))
        
        # 转换文件
        success_count = 0
        for file_path in sample_files:
            print(f"\n📄 转换文件: {file_path.name}")
            try:
                if converter.convert_file(file_path):
                    success_count += 1
                    pdf_file = output_dir / f"{file_path.stem}.pdf"
                    if pdf_file.exists():
                        file_size = pdf_file.stat().st_size / 1024  # KB
                        print(f"✅ 转换成功: {pdf_file.name} ({file_size:.1f} KB)")
                    else:
                        print(f"⚠️  转换命令成功但未找到PDF文件")
                else:
                    print(f"❌ 转换失败: {file_path.name}")
            except Exception as e:
                print(f"❌ 转换异常: {e}")
        
        # 总结
        print(f"\n📊 转换结果:")
        print(f"   总文件数: {len(sample_files)}")
        print(f"   成功转换: {success_count}")
        print(f"   失败数量: {len(sample_files) - success_count}")
        
        if success_count > 0:
            print(f"\n📁 输出目录: {output_dir}")
            if output_dir.exists():
                pdf_files = list(output_dir.glob("*.pdf"))
                print(f"   生成PDF文件: {len(pdf_files)} 个")
                for pdf_file in pdf_files:
                    size_kb = pdf_file.stat().st_size / 1024
                    print(f"   - {pdf_file.name} ({size_kb:.1f} KB)")
        
        print(f"\n🎉 演示完成！")
        
        # 提示用户
        if success_count > 0:
            print(f"\n💡 提示: PDF文件已生成在临时目录中")
            print(f"   如需保留，请在程序结束前复制文件")
            input("按回车键继续...")


if __name__ == "__main__":
    demo_conversion()
